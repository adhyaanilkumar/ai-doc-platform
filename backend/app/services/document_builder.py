import io
from typing import Iterable
from fastapi import HTTPException, status
from fastapi.responses import StreamingResponse
from docx import Document
from pptx import Presentation


def export_docx(project, sections: Iterable):
    if project.document_type != "docx":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Project is not configured as a Word document.",
        )

    doc = Document()
    doc.add_heading(project.title, 0)
    for section in sorted(sections, key=lambda s: s.order_index):
        doc.add_heading(section.title, level=1)
        content = section.content or ""
        for paragraph in content.split("\n\n"):
            doc.add_paragraph(paragraph.strip())

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{project.title}.docx"'},
    )


def export_pptx(project, sections: Iterable):
    if project.document_type != "pptx":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Project is not configured as a PowerPoint document.",
        )

    presentation = Presentation()
    for section in sorted(sections, key=lambda s: s.order_index):
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = section.title
        body_shape = slide.shapes.placeholders[1]
        body_shape.text = section.content or ""

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{project.title}.pptx"'},
    )

